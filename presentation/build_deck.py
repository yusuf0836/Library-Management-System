"""
Generates a professional presentation deck for the Library Management System project.

Usage:  python3 presentation/build_deck.py
Output: presentation/Library-Management-System-Presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

# ----------------------------------------------------------------------------
# Design tokens
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x0E, 0x1B, 0x33)
NAVY_SOFT = RGBColor(0x1B, 0x2C, 0x4F)
TEAL = RGBColor(0x1F, 0xB6, 0xA6)
AMBER = RGBColor(0xF4, 0xB4, 0x3F)
CORAL = RGBColor(0xE9, 0x6A, 0x5E)
VIOLET = RGBColor(0x7C, 0x6C, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF6, 0xF8, 0xFB)
GREY_TXT = RGBColor(0x4A, 0x55, 0x68)
GREY_LINE = RGBColor(0xDD, 0xE3, 0xEC)
DARK_TXT = RGBColor(0x12, 0x1C, 0x2E)

FONT_H = "Trebuchet MS"
FONT_B = "Segoe UI"

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Primitives
# ----------------------------------------------------------------------------
def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         shadow=False, adj=None):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    if not shadow:
        s.shadow.inherit = False
    if adj is not None:
        try:
            s.adjustments[0] = adj
        except Exception:
            pass
    s.text_frame.word_wrap = True
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, space_after=0):
    """runs: list of (text, size, bold, color, font) or list of paragraphs (list of runs)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paragraphs = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for item in para:
            t, size, bold, color = item[0], item[1], item[2], item[3]
            font_name = item[4] if len(item) > 4 else FONT_B
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font_name
    return tb


def fit_text(shape, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for item in para:
            t, size, bold, color = item[0], item[1], item[2], item[3]
            font_name = item[4] if len(item) > 4 else FONT_B
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font_name


SLIDE_NO = {"n": 0}


def base_slide(title, kicker=None, dark=False):
    """Standard content slide with header band."""
    SLIDE_NO["n"] += 1
    s = prs.slides.add_slide(BLANK)
    bg = NAVY if dark else OFFWHITE
    rect(s, 0, 0, SW, SH, fill=bg)

    # accent bar top-left
    rect(s, Inches(0.62), Inches(0.52), Inches(0.09), Inches(0.62), fill=TEAL)

    title_color = WHITE if dark else DARK_TXT
    kick_color = TEAL if dark else TEAL
    y = Inches(0.50)
    if kicker:
        text(s, Inches(0.88), y, Inches(10.5), Inches(0.24),
             [(kicker.upper(), 11, True, kick_color)])
        y = Inches(0.74)
    text(s, Inches(0.88), y, Inches(10.8), Inches(0.5),
         [(title, 28, True, title_color, FONT_H)])

    # footer
    rect(s, Inches(0.62), Inches(6.86), Inches(12.11), Pt(0.75),
         fill=(NAVY_SOFT if dark else GREY_LINE))
    foot_c = RGBColor(0x8C, 0x9A, 0xB0)
    text(s, Inches(0.62), Inches(6.98), Inches(7.0), Inches(0.25),
         [("Library Management System  ·  Laravel + MySQL", 9.5, False, foot_c)])
    text(s, Inches(10.0), Inches(6.98), Inches(2.73), Inches(0.25),
         [(f"{SLIDE_NO['n']:02d}", 9.5, True, foot_c)], align=PP_ALIGN.RIGHT)
    return s


def card(slide, x, y, w, h, accent, heading, body, num=None, hsize=14, bsize=10.5):
    c = rect(slide, x, y, w, h, fill=WHITE, line=GREY_LINE, line_w=0.75)
    rect(slide, x, y, Inches(0.055), h, fill=accent)
    tf = c.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.26)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.17)
    tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]
    if num:
        r = p.add_run()
        r.text = num + "   "
        r.font.size = Pt(hsize)
        r.font.bold = True
        r.font.color.rgb = accent
        r.font.name = FONT_H
    r = p.add_run()
    r.text = heading
    r.font.size = Pt(hsize)
    r.font.bold = True
    r.font.color.rgb = DARK_TXT
    r.font.name = FONT_H
    if body:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        p2.line_spacing = 1.18
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(bsize)
        r2.font.color.rgb = GREY_TXT
        r2.font.name = FONT_B
    return c


def bullets(slide, x, y, w, items, size=12.5, gap=0.42, dot=TEAL, bold_lead=True):
    for i, it in enumerate(items):
        yy = y + Inches(gap * i)
        rect(slide, x, yy + Inches(0.085), Inches(0.11), Inches(0.11),
             fill=dot, shape=MSO_SHAPE.OVAL)
        if isinstance(it, tuple):
            runs = [(it[0], size, bold_lead, DARK_TXT), (it[1], size, False, GREY_TXT)]
        else:
            runs = [(it, size, False, GREY_TXT)]
        text(slide, x + Inches(0.26), yy, w, Inches(0.34), runs, line_spacing=1.15)


def table(slide, x, y, w, headers, rows, col_ratios=None, row_h=0.36, head_h=0.42,
          fsize=11, head_bg=NAVY, zebra=RGBColor(0xF2, 0xF5, 0xFA)):
    ncols = len(headers)
    if col_ratios is None:
        col_ratios = [1] * ncols
    total = sum(col_ratios)
    widths = [Emu(int(w * r / total)) for r in col_ratios]

    # header
    cx = x
    for i, hcell in enumerate(headers):
        cell = rect(slide, cx, y, widths[i], Inches(head_h), fill=head_bg)
        fit_text(cell, [[(hcell, fsize, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
        cx += widths[i]

    for ri, row in enumerate(rows):
        ry = y + Inches(head_h) + Inches(row_h * ri)
        cx = x
        bg = WHITE if ri % 2 == 0 else zebra
        for ci, val in enumerate(row):
            cell = rect(slide, cx, ry, widths[ci], Inches(row_h), fill=bg,
                        line=GREY_LINE, line_w=0.5)
            bold = ci == 0
            fit_text(cell, [[(val, fsize - 0.5, bold, DARK_TXT if bold else GREY_TXT)]],
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += widths[ci]


def stat(slide, x, y, w, h, value, label, accent):
    box = rect(slide, x, y, w, h, fill=WHITE, line=GREY_LINE, line_w=0.75)
    rect(slide, x, y, w, Inches(0.07), fill=accent)
    fit_text(box, [
        [(value, 30, True, accent, FONT_H)],
        [(label, 10.5, False, GREY_TXT)],
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)


def pill(slide, x, y, w, h, label, fill, txt=WHITE, size=10.5, bold=True):
    p = rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.28)
    fit_text(p, [[(label, size, bold, txt)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    return p


def arrow(slide, x, y, w, h=Inches(0.16), color=RGBColor(0xA8, 0xB4, 0xC6)):
    a = rect(slide, x, y, w, h, fill=color, shape=MSO_SHAPE.RIGHT_ARROW)
    return a


# ============================================================================
# 1. TITLE
# ============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
# decorative geometry
rect(s, Inches(9.1), Inches(-1.4), Inches(6.6), Inches(6.6), fill=NAVY_SOFT,
     shape=MSO_SHAPE.OVAL)
rect(s, Inches(10.9), Inches(3.6), Inches(4.2), Inches(4.2), fill=RGBColor(0x15, 0x24, 0x43),
     shape=MSO_SHAPE.OVAL)
rect(s, Inches(0), Inches(0), Inches(0.16), SH, fill=TEAL)

text(s, Inches(1.0), Inches(1.45), Inches(9.0), Inches(0.3),
     [("ACADEMIC PROJECT PRESENTATION  ·  2026", 12, True, TEAL)])
text(s, Inches(1.0), Inches(1.95), Inches(9.4), Inches(1.9),
     [[("Library Management", 47, True, WHITE, FONT_H)],
      [("System", 47, True, TEAL, FONT_H)]], line_spacing=1.02)
rect(s, Inches(1.0), Inches(4.05), Inches(1.5), Pt(3), fill=AMBER)
text(s, Inches(1.0), Inches(4.35), Inches(8.4), Inches(0.9),
     [("A full-stack web application that digitizes cataloguing, membership, "
       "circulation, fines and reporting for an academic library.",
       15, False, RGBColor(0xB6, 0xC4, 0xD8))], line_spacing=1.3)

for i, (lbl, col) in enumerate([("Laravel 13", TEAL), ("PHP 8.3", AMBER),
                                ("MySQL", VIOLET), ("Bootstrap 5", CORAL)]):
    pill(s, Inches(1.0 + i * 1.72), Inches(5.5), Inches(1.55), Inches(0.4), lbl,
         RGBColor(0x1B, 0x2C, 0x4F), col, size=10.5)

rect(s, Inches(1.0), Inches(6.35), Inches(11.3), Pt(0.75), fill=RGBColor(0x27, 0x3A, 0x5E))
text(s, Inches(1.0), Inches(6.55), Inches(6.0), Inches(0.3),
     [("Presented by  ", 11, False, RGBColor(0x7E, 0x8E, 0xA8)),
      ("Yusuf", 11, True, WHITE)])
text(s, Inches(7.0), Inches(6.55), Inches(5.3), Inches(0.3),
     [("github.com/yusuf0836/Library-Management-System", 11, False,
       RGBColor(0x7E, 0x8E, 0xA8))], align=PP_ALIGN.RIGHT)

# ============================================================================
# 2. AGENDA
# ============================================================================
s = base_slide("Agenda", "What this presentation covers")
items = [
    ("01", "Problem Statement", "Why manual library operations break down"),
    ("02", "Objectives & Scope", "What the system sets out to solve"),
    ("03", "Technology Stack", "Laravel, MySQL, Blade, Chart.js"),
    ("04", "System Architecture", "MVC request lifecycle and layers"),
    ("05", "Database Design", "11 tables, relationships, integrity"),
    ("06", "Roles & Security", "Admin, Librarian, Member access control"),
    ("07", "Core Modules", "Catalogue, membership, circulation, fines"),
    ("08", "Reports & Dashboard", "Live KPIs, charts and CSV export"),
    ("09", "Testing & Results", "Feature tests and outcomes"),
    ("10", "Limitations & Future Work", "Roadmap beyond version 1"),
]
for i, (n, h, b) in enumerate(items):
    col = i % 2
    row = i // 2
    x = Inches(0.62 + col * 6.15)
    y = Inches(1.62 + row * 1.05)
    card(s, x, y, Inches(5.9), Inches(0.9), [TEAL, AMBER, VIOLET, CORAL][i % 4],
         h, b, num=n, hsize=13, bsize=10)

# ============================================================================
# 3. PROBLEM STATEMENT
# ============================================================================
s = base_slide("The Problem with Manual Library Operations", "01 · Problem statement")
probs = [
    (TEAL, "Registers don't scale", "Ledger books and spreadsheets cannot answer "
     "\"who has this copy right now?\" without a manual search."),
    (AMBER, "No real-time availability", "Students queue at the desk only to learn "
     "every copy of a title is already issued."),
    (CORAL, "Overdue tracking fails", "Due dates are checked by eye, so late returns "
     "and fines are inconsistently enforced."),
    (VIOLET, "No usable reporting", "Circulation trends, stock status and outstanding "
     "dues take hours to compile by hand."),
]
for i, (c, h, b) in enumerate(probs):
    card(s, Inches(0.62 + (i % 2) * 6.15), Inches(1.7 + (i // 2) * 1.55),
         Inches(5.9), Inches(1.35), c, h, b)

band = rect(s, Inches(0.62), Inches(4.95), Inches(12.11), Inches(1.05), fill=NAVY)
fit_text(band, [[("Consequence:  ", 13, True, AMBER),
                 ("staff time is consumed by clerical work, book loss goes undetected, "
                  "and library policy cannot be enforced consistently or audited.",
                  13, False, RGBColor(0xD5, 0xDE, 0xEA))]],
         anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# 4. OBJECTIVES
# ============================================================================
s = base_slide("Project Objectives", "02 · Goals")
objs = [
    "Digitize the full catalogue — categories, authors, publishers, books and physical copies.",
    "Track every physical copy individually by accession number, shelf location and status.",
    "Automate the issue and return workflow with enforced due dates.",
    "Generate overdue fines automatically from configurable library rules.",
    "Enforce role-based access for Admin, Librarian and Member users.",
    "Provide a live dashboard with statistics, charts and exportable reports.",
]
bullets(s, Inches(0.72), Inches(1.75), Inches(7.1), objs, size=12.5, gap=0.62)

panel = rect(s, Inches(8.35), Inches(1.62), Inches(4.38), Inches(4.9), fill=NAVY)
text(s, Inches(8.68), Inches(1.95), Inches(3.8), Inches(0.3),
     [("SCOPE BOUNDARY", 11, True, TEAL)])
rect(s, Inches(8.68), Inches(2.32), Inches(0.9), Pt(2), fill=AMBER)
in_scope = ["Catalogue & copy management", "Member lifecycle", "Issue / return circulation",
            "Fine calculation & payment", "Dashboard, reports, CSV export", "System settings"]
for i, t in enumerate(in_scope):
    text(s, Inches(8.68), Inches(2.62 + i * 0.36), Inches(3.7), Inches(0.3),
         [("✓  ", 12, True, TEAL), (t, 11, False, RGBColor(0xC8, 0xD3, 0xE2))])
text(s, Inches(8.68), Inches(4.95), Inches(3.7), Inches(0.3),
     [("OUT OF SCOPE (v1)", 10.5, True, CORAL)])
for i, t in enumerate(["Online reservations & payment gateway", "Barcode / RFID hardware",
                       "Email & SMS notifications"]):
    text(s, Inches(8.68), Inches(5.3 + i * 0.34), Inches(3.7), Inches(0.3),
         [("×  ", 12, True, CORAL), (t, 10.5, False, RGBColor(0x9E, 0xAC, 0xC2))])

# ============================================================================
# 5. TECHNOLOGY STACK
# ============================================================================
s = base_slide("Technology Stack", "03 · Tools & frameworks")
rows = [
    ("Backend", "PHP 8.3 · Laravel 13", "Routing, controllers, validation, middleware"),
    ("Database", "MySQL", "Relational storage with foreign-key integrity"),
    ("ORM", "Laravel Eloquent", "Models, relationships, eager loading"),
    ("Frontend", "Blade · Bootstrap 5 · JS", "Server-rendered responsive interface"),
    ("Charts", "Chart.js", "Copy-status doughnut and monthly issue trend"),
    ("Auth", "Laravel session auth", "Login, hashing, CSRF, role middleware"),
    ("Build", "Vite · npm", "Asset bundling for CSS and JavaScript"),
    ("Tooling", "Git · GitHub · XAMPP · PHPUnit", "Version control, local server, testing"),
]
table(s, Inches(0.62), Inches(1.68), Inches(12.11),
      ["Layer", "Technology", "Role in the system"],
      rows, col_ratios=[2.0, 3.2, 5.2], row_h=0.42, head_h=0.46, fsize=11.5)

note = rect(s, Inches(0.62), Inches(5.55), Inches(12.11), Inches(0.9), fill=WHITE,
            line=GREY_LINE, line_w=0.75)
rect(s, Inches(0.62), Inches(5.55), Inches(0.055), Inches(0.9), fill=TEAL)
fit_text(note, [[("Why Laravel?  ", 12, True, DARK_TXT),
                 ("Eloquent ORM, migrations, request validation, middleware and Blade "
                  "templating remove large amounts of boilerplate, so development time "
                  "goes into library logic instead of plumbing.", 11.5, False, GREY_TXT)]],
         anchor=MSO_ANCHOR.MIDDLE)

# ============================================================================
# 6. ARCHITECTURE
# ============================================================================
s = base_slide("System Architecture — MVC Request Lifecycle", "04 · Architecture")

lanes = [
    ("Browser", "Blade views\nBootstrap 5 UI", TEAL),
    ("Routes", "routes/web.php\nnamed routes", AMBER),
    ("Middleware", "auth · guest\nrole:admin,librarian", CORAL),
    ("Controllers", "15 controllers\nvalidation + logic", VIOLET),
    ("Models", "10 Eloquent models\nrelationships", TEAL),
    ("MySQL", "11 tables\nFK constraints", AMBER),
]
x0 = Inches(0.62)
bw = Inches(1.79)
gap = Inches(0.26)
for i, (title_, sub, col) in enumerate(lanes):
    x = Emu(int(x0 + i * (bw + gap)))
    b = rect(s, x, Inches(1.85), bw, Inches(1.55), fill=WHITE, line=GREY_LINE, line_w=0.75)
    rect(s, x, Inches(1.85), bw, Inches(0.09), fill=col)
    fit_text(b, [[(title_, 12.5, True, DARK_TXT, FONT_H)],
                 [(sub.replace("\n", "  "), 9.5, False, GREY_TXT)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.1)
    if i < len(lanes) - 1:
        arrow(s, Emu(int(x + bw + Inches(0.04))), Inches(2.54), Inches(0.18), Inches(0.18))

text(s, Inches(0.62), Inches(3.58), Inches(12.11), Inches(0.3),
     [("Request flows left to right; the rendered HTML response flows back.",
       10.5, False, RGBColor(0x8C, 0x9A, 0xB0))], align=PP_ALIGN.CENTER)

# lower detail cards
det = [
    (TEAL, "Separation of concerns",
     "Controllers hold request handling and business rules; models own data access and "
     "relationships; Blade owns presentation only."),
    (AMBER, "Guarded routes",
     "Route groups apply auth and role middleware, so authorization is declared once "
     "per group rather than repeated per action."),
    (VIOLET, "Transactional integrity",
     "Issue, return and fine-payment operations run inside DB transactions with row "
     "locking to prevent double-issuing a copy."),
]
for i, (c, h, b) in enumerate(det):
    card(s, Inches(0.62 + i * 4.11), Inches(4.1), Inches(3.87), Inches(1.55), c, h, b,
         hsize=12.5, bsize=10)

# ============================================================================
# 7. DATABASE DESIGN
# ============================================================================
s = base_slide("Database Design — 11 Tables", "05 · Data model")
db_rows = [
    ("users", "Login credentials, role, profile photo"),
    ("members", "Member profile, member code, department, active flag"),
    ("categories", "Book categories"),
    ("authors", "Author records and biographies"),
    ("publishers", "Publisher information"),
    ("books", "Title, ISBN, edition, description, cover image"),
    ("author_book", "Pivot table: many-to-many books ↔ authors"),
    ("book_copies", "Physical copy, accession number, shelf, status"),
    ("book_issues", "Issue, due, return and circulation records"),
    ("fines", "Overdue amount, paid amount, payment status"),
    ("settings", "Library rules: borrowing period, fine rate, contact"),
]
table(s, Inches(0.62), Inches(1.68), Inches(7.55), ["Table", "Purpose"],
      db_rows, col_ratios=[2.1, 5.4], row_h=0.365, head_h=0.4, fsize=11)

p = rect(s, Inches(8.5), Inches(1.68), Inches(4.23), Inches(4.5), fill=NAVY)
text(s, Inches(8.82), Inches(1.98), Inches(3.6), Inches(0.3),
     [("KEY RELATIONSHIPS", 11, True, TEAL)])
rect(s, Inches(8.82), Inches(2.35), Inches(0.9), Pt(2), fill=AMBER)
rels = [
    ("Book → BookCopy", "one-to-many"),
    ("Book ↔ Author", "many-to-many"),
    ("Book → Category / Publisher", "many-to-one"),
    ("User → Member", "one-to-one"),
    ("Member → BookIssue", "one-to-many"),
    ("BookCopy → BookIssue", "one-to-many"),
    ("BookIssue → Fine", "one-to-one"),
]
for i, (a, b) in enumerate(rels):
    yy = Inches(2.68 + i * 0.49)
    text(s, Inches(8.82), yy, Inches(3.6), Inches(0.24),
         [(a, 11, True, WHITE)])
    text(s, Inches(8.82), yy + Inches(0.21), Inches(3.6), Inches(0.22),
         [(b, 9.5, False, TEAL)])

# ============================================================================
# 8. ER OVERVIEW (visual)
# ============================================================================
s = base_slide("Entity Relationship Overview", "05 · Data model")


def entity(x, y, w, h, name, fields, col):
    b = rect(s, x, y, w, h, fill=WHITE, line=GREY_LINE, line_w=0.75)
    head = rect(s, x, y, w, Inches(0.34), fill=col)
    fit_text(head, [[(name, 11, True, WHITE)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.16), y + Inches(0.44), w - Inches(0.28), h - Inches(0.5),
         [[(f, 9.5, False, GREY_TXT)] for f in fields], line_spacing=1.25)
    return b


entity(Inches(0.62), Inches(1.72), Inches(2.5), Inches(1.5), "CATEGORY",
       ["id", "name", "description"], VIOLET)
entity(Inches(0.62), Inches(3.52), Inches(2.5), Inches(1.5), "PUBLISHER",
       ["id", "name", "email, phone"], VIOLET)
entity(Inches(0.62), Inches(5.32), Inches(2.5), Inches(1.2), "AUTHOR",
       ["id", "name, bio"], VIOLET)

entity(Inches(3.85), Inches(2.55), Inches(2.6), Inches(2.1), "BOOK",
       ["id", "title, isbn", "edition", "category_id (FK)", "publisher_id (FK)",
        "cover_image"], TEAL)

entity(Inches(7.05), Inches(2.55), Inches(2.6), Inches(2.1), "BOOK_COPY",
       ["id", "book_id (FK)", "accession_number", "shelf_location",
        "status"], TEAL)

entity(Inches(10.25), Inches(1.72), Inches(2.48), Inches(1.95), "BOOK_ISSUE",
       ["id", "member_id (FK)", "book_copy_id (FK)", "issued_at, due_at",
        "returned_at, status"], AMBER)
entity(Inches(10.25), Inches(4.0), Inches(2.48), Inches(1.4), "FINE",
       ["id", "book_issue_id (FK)", "amount, paid_amount, status"], CORAL)

entity(Inches(7.05), Inches(5.05), Inches(2.6), Inches(1.45), "MEMBER",
       ["id", "user_id (FK)", "member_code", "is_active"], AMBER)
entity(Inches(3.85), Inches(5.05), Inches(2.6), Inches(1.45), "USER",
       ["id", "name, email", "password", "role"], NAVY_SOFT)


def connect(x1, y1, x2, y2, label=None):
    ln = s.shapes.add_connector(1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    ln.line.color.rgb = RGBColor(0xA8, 0xB4, 0xC6)
    ln.line.width = Pt(1.25)
    if label:
        text(s, Emu(int((x1 + x2) / 2 - Inches(0.35))),
             Emu(int((y1 + y2) / 2 - Inches(0.22))), Inches(0.7), Inches(0.2),
             [(label, 8.5, True, RGBColor(0x8C, 0x9A, 0xB0))], align=PP_ALIGN.CENTER)


connect(Inches(3.12), Inches(2.47), Inches(3.85), Inches(3.15), "1:N")
connect(Inches(3.12), Inches(4.27), Inches(3.85), Inches(3.9), "1:N")
connect(Inches(3.12), Inches(5.92), Inches(3.85), Inches(4.5), "M:N")
connect(Inches(6.45), Inches(3.6), Inches(7.05), Inches(3.6), "1:N")
connect(Inches(9.65), Inches(3.0), Inches(10.25), Inches(2.7), "1:N")
connect(Inches(9.65), Inches(5.5), Inches(10.25), Inches(3.3), "1:N")
connect(Inches(11.49), Inches(3.67), Inches(11.49), Inches(4.0), "1:1")
connect(Inches(6.45), Inches(5.78), Inches(7.05), Inches(5.78), "1:1")

# ============================================================================
# 9. ROLES & SECURITY
# ============================================================================
s = base_slide("Role-Based Access Control", "06 · Security")
roles = [
    (CORAL, "ADMIN", "Full system authority",
     ["Everything a Librarian can do", "System settings & library rules",
      "Fine rate and borrowing period", "User & role administration"]),
    (TEAL, "LIBRARIAN", "Daily desk operations",
     ["Books, copies, authors, categories", "Member records & activation",
      "Issue and return circulation", "Fines, reports and dashboard"]),
    (AMBER, "MEMBER", "Self-service only",
     ["Personal dashboard", "Own borrowing history", "Own due dates",
      "Own fine status"]),
]
for i, (col, name, sub, perms) in enumerate(roles):
    x = Inches(0.62 + i * 4.11)
    b = rect(s, x, Inches(1.68), Inches(3.87), Inches(3.2), fill=WHITE,
             line=GREY_LINE, line_w=0.75)
    head = rect(s, x, Inches(1.68), Inches(3.87), Inches(0.62), fill=col)
    fit_text(head, [[(name, 15, True, WHITE, FONT_H)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.24), Inches(2.46), Inches(3.4), Inches(0.24),
         [(sub, 10.5, True, col)])
    for j, p_ in enumerate(perms):
        text(s, x + Inches(0.24), Inches(2.82 + j * 0.44), Inches(3.4), Inches(0.4),
             [("•  ", 11, True, col), (p_, 10.5, False, GREY_TXT)], line_spacing=1.1)

sec = rect(s, Inches(0.62), Inches(5.08), Inches(12.11), Inches(1.4), fill=NAVY)
text(s, Inches(0.95), Inches(5.3), Inches(11.4), Inches(0.28),
     [("SECURITY MEASURES", 11, True, TEAL)])
measures = ["Hashed passwords", "CSRF protection on all forms",
            "Server-side validation", "Route-group middleware guards",
            "Inactive members blocked at login", "403 on unauthorized access"]
for i, m in enumerate(measures):
    text(s, Inches(0.95 + (i % 3) * 3.95), Inches(5.7 + (i // 3) * 0.34),
         Inches(3.8), Inches(0.28),
         [("✓  ", 11, True, AMBER), (m, 11, False, RGBColor(0xC8, 0xD3, 0xE2))])

# ============================================================================
# 10. MODULE MAP
# ============================================================================
s = base_slide("Core Modules", "07 · Functionality")
mods = [
    (TEAL, "Catalogue", "Categories, authors, publishers, books with cover upload, "
     "search and filtering, multi-author support."),
    (AMBER, "Copy Inventory", "Per-copy accession number, shelf location and status: "
     "available, issued, reserved, lost, damaged."),
    (VIOLET, "Membership", "Member code, department, contact details, linked login "
     "account, active/inactive control."),
    (CORAL, "Circulation", "Issue to active members only, due-date tracking, return "
     "processing, automatic copy-status update."),
    (TEAL, "Fines", "Automatic overdue fine, unpaid/partial/paid/waived states, "
     "payment recording, outstanding totals."),
    (AMBER, "Reports", "Overdue and circulation reports, date and status filters, "
     "CSV export for offline records."),
    (VIOLET, "Dashboard", "Role-aware KPI cards, copy-status doughnut chart, "
     "six-month issue trend, recent activity."),
    (CORAL, "Settings", "Admin-only library name and contact, default borrowing "
     "days, fine rate per overdue day."),
]
for i, (c, h, b) in enumerate(mods):
    card(s, Inches(0.62 + (i % 4) * 3.09), Inches(1.72 + (i // 4) * 2.4),
         Inches(2.87), Inches(2.2), c, h, b, hsize=13, bsize=9.8)

# ============================================================================
# 11. CIRCULATION WORKFLOW
# ============================================================================
s = base_slide("Circulation Workflow — Issue to Return", "07 · Key process")

steps = [
    ("1", "Select member", "Only active members appear in the list"),
    ("2", "Pick available copy", "Copy is locked for update inside a transaction"),
    ("3", "Set dates", "Due date defaults from the borrowing-period setting"),
    ("4", "Copy → issued", "Status flips and the issue record is created"),
    ("5", "Return recorded", "Return date validated against the issue date"),
    ("6", "Fine if late", "overdue days × fine rate → unpaid fine record"),
]
for i, (n, h, b) in enumerate(steps):
    x = Inches(0.62 + i * 2.06)
    circ = rect(s, x + Inches(0.72), Inches(1.78), Inches(0.56), Inches(0.56),
                fill=NAVY, shape=MSO_SHAPE.OVAL)
    fit_text(circ, [[(n, 14, True, TEAL, FONT_H)]], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    b_ = rect(s, x, Inches(2.55), Inches(1.94), Inches(1.5), fill=WHITE,
              line=GREY_LINE, line_w=0.75)
    fit_text(b_, [[(h, 11.5, True, DARK_TXT, FONT_H)], [(b, 9, False, GREY_TXT)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.12)
    if i < 5:
        arrow(s, x + Inches(1.96), Inches(2.03), Inches(0.14), Inches(0.14))

hl = rect(s, Inches(0.62), Inches(4.35), Inches(5.95), Inches(2.1), fill=NAVY)
text(s, Inches(0.95), Inches(4.6), Inches(5.3), Inches(0.28),
     [("BUSINESS RULES ENFORCED IN CODE", 11, True, TEAL)])
for i, t in enumerate(["A copy that is not 'available' can never be issued",
                       "Inactive members are rejected at issue time",
                       "Return date must be on or after the issue date",
                       "A returned book cannot be returned twice"]):
    text(s, Inches(0.95), Inches(5.0 + i * 0.34), Inches(5.3), Inches(0.28),
         [("→  ", 11, True, AMBER), (t, 10.5, False, RGBColor(0xC8, 0xD3, 0xE2))])

f = rect(s, Inches(6.78), Inches(4.35), Inches(5.95), Inches(2.1), fill=WHITE,
         line=GREY_LINE, line_w=0.75)
rect(s, Inches(6.78), Inches(4.35), Inches(0.055), Inches(2.1), fill=CORAL)
text(s, Inches(7.1), Inches(4.6), Inches(5.3), Inches(0.28),
     [("FINE CALCULATION", 11, True, CORAL)])
text(s, Inches(7.1), Inches(4.98), Inches(5.4), Inches(0.5),
     [("fine = overdue_days × fine_per_day", 17, True, DARK_TXT, FONT_H)])
text(s, Inches(7.1), Inches(5.52), Inches(5.4), Inches(0.85),
     [("The rate is read from the settings table, so an administrator can change "
       "library policy without touching code. Fines are written in the same "
       "transaction as the return.", 10.5, False, GREY_TXT)], line_spacing=1.2)

# ============================================================================
# 12. DASHBOARD & REPORTS
# ============================================================================
s = base_slide("Dashboard & Reporting", "08 · Insight")
kpis = [("Total Books", TEAL), ("Available Copies", TEAL), ("Issued Copies", AMBER),
        ("Active Members", VIOLET), ("Overdue Issues", CORAL), ("Outstanding Fines", CORAL)]
for i, (lbl, c) in enumerate(kpis):
    x = Inches(0.62 + i * 2.06)
    b = rect(s, x, Inches(1.7), Inches(1.94), Inches(1.05), fill=WHITE,
             line=GREY_LINE, line_w=0.75)
    rect(s, x, Inches(1.7), Inches(1.94), Inches(0.07), fill=c)
    fit_text(b, [[("KPI", 16, True, c, FONT_H)], [(lbl, 9.5, False, GREY_TXT)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)

feat = [
    (TEAL, "Copy-status chart", "Chart.js doughnut across available, issued, reserved, "
     "lost and damaged copies."),
    (AMBER, "Six-month issue trend", "Monthly issue counts aggregated with grouped SQL "
     "and rendered as a trend chart."),
    (VIOLET, "Overdue report", "Every unreturned copy past its due date, ordered by "
     "how long it is overdue."),
    (CORAL, "Circulation + CSV", "Date-range and status filters with a streamed CSV "
     "download for records."),
]
for i, (c, h, b) in enumerate(feat):
    card(s, Inches(0.62 + i * 3.09), Inches(3.0), Inches(2.87), Inches(1.75), c, h, b,
         hsize=12.5, bsize=9.8)

band = rect(s, Inches(0.62), Inches(5.05), Inches(12.11), Inches(1.42), fill=NAVY)
text(s, Inches(0.95), Inches(5.28), Inches(11.4), Inches(0.28),
     [("ROLE-AWARE DASHBOARD", 11, True, TEAL)])
text(s, Inches(0.95), Inches(5.66), Inches(11.4), Inches(0.7),
     [("The same route serves three different views. Admins and Librarians see "
       "library-wide statistics, charts and recent circulation activity; a Member sees "
       "only their own active borrowings, overdue count and unpaid fine balance — "
       "computed from their linked member record.",
       12, False, RGBColor(0xC8, 0xD3, 0xE2))], line_spacing=1.25)

# ============================================================================
# 13. IMPLEMENTATION HIGHLIGHTS
# ============================================================================
s = base_slide("Implementation Highlights", "08 · Engineering quality")
hi = [
    (TEAL, "Race-safe issuing",
     "The issue routine opens a database transaction and calls lockForUpdate() on the "
     "book copy, so two librarians cannot issue the same physical copy at once."),
    (AMBER, "Eager loading everywhere",
     "Listing screens load nested relations (member.user, copy.book, fine) in one "
     "query set, avoiding N+1 query problems on paginated tables."),
    (VIOLET, "Filters that survive paging",
     "Search and status filters use conditional when() clauses with withQueryString(), "
     "so filters persist across pages and export links."),
    (CORAL, "Streamed CSV export",
     "Circulation exports stream rows straight to the response instead of building the "
     "whole file in memory."),
    (TEAL, "Configuration, not hard-coding",
     "Borrowing period and fine rate come from the settings table with sensible "
     "fallbacks, so policy changes need no redeploy."),
    (AMBER, "Consistent validation",
     "Every write action validates input server-side before it reaches the database, "
     "with user-facing error and success feedback."),
]
for i, (c, h, b) in enumerate(hi):
    card(s, Inches(0.62 + (i % 2) * 6.15), Inches(1.72 + (i // 2) * 1.63),
         Inches(5.9), Inches(1.45), c, h, b, hsize=13, bsize=10)

# ============================================================================
# 14. TESTING
# ============================================================================
s = base_slide("Testing & Quality Assurance", "09 · Validation")
text(s, Inches(0.72), Inches(1.72), Inches(6.0), Inches(0.3),
     [("AUTOMATED TESTS (PHPUnit)", 11, True, TEAL)])
table(s, Inches(0.62), Inches(2.08), Inches(6.1), ["Test suite", "Verifies"],
      [("AuthenticationTest", "Login, logout, invalid credentials"),
       ("RoleAccessTest", "Admin / Librarian / Member route guards"),
       ("Feature tests", "HTTP responses and redirects"),
       ("Unit tests", "Isolated helper behaviour")],
      col_ratios=[2.3, 3.6], row_h=0.42, head_h=0.44, fsize=11)

text(s, Inches(7.0), Inches(1.72), Inches(5.7), Inches(0.3),
     [("MANUAL TEST SCENARIOS", 11, True, AMBER)])
table(s, Inches(6.9), Inches(2.08), Inches(5.83), ["Scenario", "Expected"],
      [("Issue an unavailable copy", "Blocked with error"),
       ("Issue to inactive member", "Blocked"),
       ("Return after due date", "Fine auto-created"),
       ("Overpay a fine", "Validation rejects")],
      col_ratios=[2.9, 2.6], row_h=0.42, head_h=0.44, fsize=11, head_bg=NAVY_SOFT)

res = [("11", "Database tables", TEAL), ("15", "Controllers", AMBER),
       ("10", "Eloquent models", VIOLET), ("3", "User roles", CORAL),
       ("8", "Functional modules", TEAL)]
for i, (v, l, c) in enumerate(res):
    stat(s, Inches(0.62 + i * 2.47), Inches(4.55), Inches(2.28), Inches(1.35), v, l, c)

text(s, Inches(0.62), Inches(6.16), Inches(12.11), Inches(0.4),
     [("Outcome: all core circulation, fine and access-control paths behave as "
       "specified under both automated and manual testing.",
       11.5, False, GREY_TXT)], align=PP_ALIGN.CENTER)

# ============================================================================
# 15. LIMITATIONS & FUTURE WORK
# ============================================================================
s = base_slide("Limitations & Future Work", "10 · Roadmap")
lim = ["No email or SMS due-date reminders",
       "No online reservation or hold queue",
       "No barcode / RFID scanning support",
       "Reports export to CSV only, not PDF",
       "Single-branch library model"]
text(s, Inches(0.72), Inches(1.72), Inches(5.4), Inches(0.3),
     [("CURRENT LIMITATIONS", 11, True, CORAL)])
bullets(s, Inches(0.72), Inches(2.15), Inches(5.2), lim, size=11.5, gap=0.46, dot=CORAL)

fut = [("Notifications — ", "automated due-date and overdue emails"),
       ("Reservations — ", "member hold queue with priority"),
       ("Barcode support — ", "scan-to-issue and scan-to-return"),
       ("REST API — ", "mobile app and OPAC kiosk clients"),
       ("Analytics — ", "most-borrowed titles, member activity ranking"),
       ("Multi-branch — ", "inter-library transfer between campuses")]
text(s, Inches(6.9), Inches(1.72), Inches(5.8), Inches(0.3),
     [("PLANNED ENHANCEMENTS", 11, True, TEAL)])
bullets(s, Inches(6.9), Inches(2.15), Inches(5.6), fut, size=11.5, gap=0.46, dot=TEAL)

band = rect(s, Inches(0.62), Inches(4.9), Inches(12.11), Inches(1.55), fill=NAVY)
text(s, Inches(0.95), Inches(5.15), Inches(11.4), Inches(0.28),
     [("VERSION 2 PRIORITY", 11, True, AMBER)])
text(s, Inches(0.95), Inches(5.52), Inches(11.4), Inches(0.75),
     [("Notifications and reservations deliver the largest gain for members with the "
       "least architectural change — both build on the existing book_issues and "
       "members tables and reuse the current settings-driven policy model.",
       12, False, RGBColor(0xC8, 0xD3, 0xE2))], line_spacing=1.25)

# ============================================================================
# 16. CONCLUSION
# ============================================================================
s = base_slide("Conclusion", "Summary")
concl = [
    (TEAL, "Objectives met", "Every objective set at the start of the project — "
     "catalogue, membership, circulation, fines, roles and reporting — is implemented "
     "and working end to end."),
    (AMBER, "Clean, extensible design", "A conventional Laravel MVC structure with a "
     "normalised schema means new modules can be added without reworking existing code."),
    (VIOLET, "Real operational value", "Desk staff get instant availability, enforced "
     "due dates and automatic fines; management gets live statistics and exportable "
     "reports."),
]
for i, (c, h, b) in enumerate(concl):
    card(s, Inches(0.62 + i * 4.11), Inches(1.75), Inches(3.87), Inches(2.15), c, h, b,
         hsize=13.5, bsize=10.5)

q = rect(s, Inches(0.62), Inches(4.2), Inches(12.11), Inches(2.25), fill=NAVY)
rect(s, Inches(0.62), Inches(4.2), Inches(12.11), Inches(0.09), fill=TEAL)
text(s, Inches(1.1), Inches(4.62), Inches(11.1), Inches(0.9),
     [("\u201CThe system replaces a manual register with an auditable, role-secured "
       "digital workflow — turning library data into decisions.\u201D",
       17, False, WHITE, FONT_H)], line_spacing=1.3)
rect(s, Inches(1.1), Inches(5.72), Inches(1.3), Pt(2), fill=AMBER)
text(s, Inches(1.1), Inches(5.92), Inches(11.1), Inches(0.3),
     [("Library Management System  ·  Laravel 13  ·  MySQL  ·  "
       "github.com/yusuf0836/Library-Management-System",
       10.5, False, RGBColor(0x8C, 0x9A, 0xB0))])

# ============================================================================
# 17. THANK YOU
# ============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, Inches(-2.0), Inches(3.2), Inches(7.0), Inches(7.0), fill=NAVY_SOFT,
     shape=MSO_SHAPE.OVAL)
rect(s, Inches(0), Inches(0), SW, Inches(0.16), fill=TEAL)
text(s, Inches(1.0), Inches(2.55), Inches(11.3), Inches(1.2),
     [("Thank You", 54, True, WHITE, FONT_H)], align=PP_ALIGN.CENTER)
rect(s, Inches(6.17), Inches(3.85), Inches(1.0), Pt(3), fill=AMBER)
text(s, Inches(1.0), Inches(4.15), Inches(11.3), Inches(0.4),
     [("Questions & Discussion", 17, False, TEAL)], align=PP_ALIGN.CENTER)
text(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.3),
     [("github.com/yusuf0836/Library-Management-System", 12, False,
       RGBColor(0x8C, 0x9A, 0xB0))], align=PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Library-Management-System-Presentation.pptx")
prs.save(out)
print("saved:", out, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
